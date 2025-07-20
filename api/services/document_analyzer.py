import os
import logging
import tempfile
from typing import Dict, List, Optional, Tuple
from pathlib import Path
import PyPDF2
import docx
import pptx
from services.gemini_utils import gemini_service
import re

logger = logging.getLogger(__name__)

class DocumentAnalyzer:
    """
    Service for analyzing documents and extracting comprehensive insights.
    Provides summary, strengths/weaknesses, questions, recommendations, and numerical analysis.
    """
    
    def __init__(self):
        self.max_content_length = 8000  # Limit content length for API calls
        self.min_content_length = 50    # Minimum content length for analysis
    
    async def analyze_document(self, file_path: str, file_extension: str, filename: str) -> Dict:
        """
        Main method to analyze a document comprehensively.
        
        Args:
            file_path: Path to the uploaded file
            file_extension: File extension (.pdf, .docx, etc.)
            filename: Original filename
            
        Returns:
            Dictionary containing comprehensive analysis
        """
        try:
            # Extract text content from document
            content = await self._extract_text_content(file_path, file_extension)
            
            if not content or len(content.strip()) < self.min_content_length:
                raise ValueError("Document content is too short or empty for analysis")
            
            # Clean and prepare content
            cleaned_content = self._clean_content(content)
            
            # Generate comprehensive analysis
            summary = await self._generate_summary(cleaned_content, filename)
            strengths_weaknesses = await self._analyze_strengths_weaknesses(cleaned_content, filename)
            questions = await self._generate_exploration_questions(cleaned_content, filename)
            recommendations = await self._generate_recommendations(cleaned_content, filename)
            numerical_analysis = await self._analyze_numerical_data(cleaned_content, filename)
            
            # Analyze document structure and type
            doc_info = self._analyze_document_structure(cleaned_content, filename)
            
            return {
                "summary": summary,
                "strengths_weaknesses": strengths_weaknesses,
                "exploration_questions": questions,
                "recommendations": recommendations,
                "numerical_analysis": numerical_analysis,
                "document_info": doc_info,
                "word_count": len(cleaned_content.split()),
                "content_preview": cleaned_content[:200] + "..." if len(cleaned_content) > 200 else cleaned_content
            }
            
        except Exception as e:
            logger.error(f"Error analyzing document: {str(e)}")
            raise Exception(f"Failed to analyze document: {str(e)}")
    
    async def _extract_text_content(self, file_path: str, file_extension: str) -> str:
        """Extract text content from various document formats."""
        content = ""
        
        try:
            if file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    
            elif file_extension == '.pdf':
                content = await self._extract_pdf_content(file_path)
                
            elif file_extension in ['.doc', '.docx']:
                content = await self._extract_word_content(file_path)
                
            elif file_extension in ['.ppt', '.pptx']:
                content = await self._extract_powerpoint_content(file_path)
                
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
                
            return content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting content from {file_extension}: {str(e)}")
            raise Exception(f"Failed to extract content from document: {str(e)}")
    
    async def _extract_pdf_content(self, file_path: str) -> str:
        """Extract text from PDF files."""
        content = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            content += f"\n--- Page {page_num + 1} ---\n"
                            content += page_text
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {page_num + 1}: {str(e)}")
                        continue
                        
            return content
            
        except Exception as e:
            raise Exception(f"Error reading PDF file: {str(e)}")
    
    async def _extract_word_content(self, file_path: str) -> str:
        """Extract text from Word documents."""
        content = ""
        try:
            doc = docx.Document(file_path)
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                content += "\n--- Table Content ---\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content += " | ".join(row_text) + "\n"
            
            return content
            
        except Exception as e:
            raise Exception(f"Error reading Word document: {str(e)}")
    
    async def _extract_powerpoint_content(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations."""
        content = ""
        try:
            presentation = pptx.Presentation(file_path)
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                content += f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n"
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                content += " | ".join(row_text) + "\n"
            
            return content
            
        except Exception as e:
            raise Exception(f"Error reading PowerPoint file: {str(e)}")
    
    def _clean_content(self, content: str) -> str:
        """Clean and normalize extracted content."""
        # Remove excessive whitespace
        content = re.sub(r'\n\s*\n', '\n\n', content)
        content = re.sub(r' +', ' ', content)
        
        # Remove page markers and other artifacts
        content = re.sub(r'--- Page \d+ ---', '', content)
        content = re.sub(r'--- Slide \d+ ---', '\n\n', content)
        content = re.sub(r'--- Table Content ---', '\n', content)
        
        # Limit content length for API processing
        if len(content) > self.max_content_length:
            content = content[:self.max_content_length] + "..."
        
        return content.strip()
    
    async def _generate_summary(self, content: str, filename: str) -> str:
        """Generate a comprehensive summary of the document."""
        prompt = f"""
Analyze the following document and provide a comprehensive summary.

Document: {filename}
Content: {content}

Please provide a clear, well-structured summary that covers:
- Main topics and objectives
- Key findings or conclusions
- Important information and insights
- Overall purpose and scope

Write in 4-6 sentences using professional, easy-to-understand language.

Summary:
"""
        
        try:
            summary = await gemini_service._generate_content(prompt)
            return summary.strip()
        except Exception as e:
            logger.error(f"Error generating summary: {str(e)}")
            return self._generate_fallback_summary(content, filename)
    
    async def _analyze_strengths_weaknesses(self, content: str, filename: str) -> Dict[str, List[str]]:
        """Analyze strengths and weaknesses of the document."""
        prompt = f"""
Analyze the following document and identify its strengths and weaknesses.

Document: {filename}
Content: {content}

Provide a balanced analysis covering:

STRENGTHS (3-5 points):
- What the document does well
- Strong arguments or evidence
- Clear explanations or structure
- Valuable insights or information

WEAKNESSES (3-5 points):
- Areas that could be improved
- Missing information or gaps
- Unclear explanations
- Potential biases or limitations

Format your response as:
STRENGTHS:
1. [First strength]
2. [Second strength]
...

WEAKNESSES:
1. [First weakness]
2. [Second weakness]
...

Analysis:
"""
        
        try:
            response = await gemini_service._generate_content(prompt)
            return self._parse_strengths_weaknesses(response)
        except Exception as e:
            logger.error(f"Error analyzing strengths/weaknesses: {str(e)}")
            return self._generate_fallback_strengths_weaknesses()
    
    async def _generate_exploration_questions(self, content: str, filename: str) -> List[str]:
        """Generate questions to aid in-depth exploration of the document."""
        prompt = f"""
Based on the following document content, generate 5-8 thought-provoking questions that would help readers explore the topic more deeply.

Document: {filename}
Content: {content}

Create questions that:
- Encourage critical thinking about the content
- Explore implications and consequences
- Challenge assumptions or conclusions
- Suggest areas for further investigation
- Connect to broader contexts or applications

Format as a numbered list:
1. [First question]
2. [Second question]
...

Questions for deeper exploration:
"""
        
        try:
            response = await gemini_service._generate_content(prompt)
            return self._parse_questions(response)
        except Exception as e:
            logger.error(f"Error generating exploration questions: {str(e)}")
            return self._generate_fallback_questions()
    
    async def _generate_recommendations(self, content: str, filename: str) -> List[str]:
        """Generate recommendations for improvement or next actions."""
        prompt = f"""
Based on the following document content, provide 5-7 actionable recommendations for improvement or next steps.

Document: {filename}
Content: {content}

Provide recommendations that address:
- How to improve the document or its implementation
- Next steps or actions to take
- Areas for further development
- Practical applications of the content
- Ways to enhance effectiveness

Format as a numbered list:
1. [First recommendation]
2. [Second recommendation]
...

Recommendations for improvement and next actions:
"""
        
        try:
            response = await gemini_service._generate_content(prompt)
            return self._parse_recommendations(response)
        except Exception as e:
            logger.error(f"Error generating recommendations: {str(e)}")
            return self._generate_fallback_recommendations()
    
    async def _analyze_numerical_data(self, content: str, filename: str) -> Dict:
        """Analyze numerical data and statistics in the document."""
        # Extract numbers and potential data patterns
        numbers = re.findall(r'\b\d+(?:\.\d+)?(?:%|\$|€|£)?\b', content)
        percentages = re.findall(r'\b\d+(?:\.\d+)?%\b', content)
        currencies = re.findall(r'[\$€£]\d+(?:\.\d+)?(?:k|K|m|M|b|B)?\b', content)
        
        if not numbers and not percentages and not currencies:
            return {
                "has_numerical_data": False,
                "summary": "No significant numerical data found in the document.",
                "key_figures": [],
                "insights": []
            }
        
        prompt = f"""
Analyze the numerical data in the following document and provide insights.

Document: {filename}
Content: {content}

Found numbers: {numbers[:20]}  # Limit to first 20 numbers
Found percentages: {percentages[:10]}
Found currency values: {currencies[:10]}

Provide analysis including:
- Key numerical findings
- Trends or patterns in the data
- Significant statistics or metrics
- Data-driven insights

If there's substantial numerical data, provide a brief analysis. If minimal data, note what's present.

Numerical Analysis:
"""
        
        try:
            if len(numbers) > 5 or len(percentages) > 2 or len(currencies) > 2:
                response = await gemini_service._generate_content(prompt)
                return {
                    "has_numerical_data": True,
                    "summary": response.strip(),
                    "key_figures": {
                        "numbers_found": len(numbers),
                        "percentages": percentages[:5],
                        "currencies": currencies[:5]
                    },
                    "insights": self._extract_numerical_insights(response)
                }
            else:
                return {
                    "has_numerical_data": True,
                    "summary": f"Limited numerical data found: {len(numbers)} numbers, {len(percentages)} percentages, {len(currencies)} currency values.",
                    "key_figures": {
                        "numbers_found": len(numbers),
                        "percentages": percentages,
                        "currencies": currencies
                    },
                    "insights": ["Document contains minimal numerical data for comprehensive analysis."]
                }
        except Exception as e:
            logger.error(f"Error analyzing numerical data: {str(e)}")
            return {
                "has_numerical_data": False,
                "summary": "Unable to analyze numerical data.",
                "key_figures": [],
                "insights": []
            }
    
    def _parse_strengths_weaknesses(self, response: str) -> Dict[str, List[str]]:
        """Parse AI response to extract strengths and weaknesses."""
        strengths = []
        weaknesses = []
        
        lines = response.split('\n')
        current_section = None
        
        for line in lines:
            line = line.strip()
            if 'STRENGTHS:' in line.upper():
                current_section = 'strengths'
                continue
            elif 'WEAKNESSES:' in line.upper():
                current_section = 'weaknesses'
                continue
            
            if line and re.match(r'^\d+\.', line):
                point = re.sub(r'^\d+\.\s*', '', line).strip()
                if current_section == 'strengths':
                    strengths.append(point)
                elif current_section == 'weaknesses':
                    weaknesses.append(point)
        
        return {
            "strengths": strengths[:5],
            "weaknesses": weaknesses[:5]
        }
    
    def _parse_questions(self, response: str) -> List[str]:
        """Parse AI response to extract questions."""
        questions = []
        lines = response.split('\n')
        
        for line in lines:
            line = line.strip()
            if re.match(r'^\d+\.', line):
                question = re.sub(r'^\d+\.\s*', '', line).strip()
                if question and question.endswith('?'):
                    questions.append(question)
        
        return questions[:8]
    
    def _parse_recommendations(self, response: str) -> List[str]:
        """Parse AI response to extract recommendations."""
        recommendations = []
        lines = response.split('\n')
        
        for line in lines:
            line = line.strip()
            if re.match(r'^\d+\.', line):
                recommendation = re.sub(r'^\d+\.\s*', '', line).strip()
                if recommendation:
                    recommendations.append(recommendation)
        
        return recommendations[:7]
    
    def _extract_numerical_insights(self, response: str) -> List[str]:
        """Extract key insights from numerical analysis."""
        insights = []
        sentences = response.split('.')
        
        for sentence in sentences:
            sentence = sentence.strip()
            if len(sentence) > 20 and any(keyword in sentence.lower() for keyword in [
                'trend', 'increase', 'decrease', 'significant', 'notable', 'key', 'important'
            ]):
                insights.append(sentence + '.')
                if len(insights) >= 5:
                    break
        
        return insights
    
    def _analyze_document_structure(self, content: str, filename: str) -> Dict:
        """Analyze document structure and determine document type."""
        word_count = len(content.split())
        
        # Determine document type based on content patterns
        doc_type = "General Document"
        if any(keyword in content.lower() for keyword in ['proposal', 'project', 'plan']):
            doc_type = "Project/Proposal Document"
        elif any(keyword in content.lower() for keyword in ['report', 'analysis', 'findings']):
            doc_type = "Report/Analysis"
        elif any(keyword in content.lower() for keyword in ['presentation', 'slide', 'overview']):
            doc_type = "Presentation"
        elif any(keyword in content.lower() for keyword in ['manual', 'guide', 'instructions']):
            doc_type = "Manual/Guide"
        elif any(keyword in content.lower() for keyword in ['research', 'study', 'methodology']):
            doc_type = "Research Document"
        
        # Estimate reading time (average 200 words per minute)
        reading_time = max(1, round(word_count / 200))
        
        return {
            "document_type": doc_type,
            "estimated_reading_time": f"{reading_time} minute{'s' if reading_time != 1 else ''}",
            "complexity": "High" if word_count > 2000 else "Medium" if word_count > 500 else "Low"
        }
    
    # Fallback methods
    def _generate_fallback_summary(self, content: str, filename: str) -> str:
        """Generate a basic summary when AI fails."""
        word_count = len(content.split())
        sentences = content.split('.')[:3]
        summary_text = '. '.join(sentence.strip() for sentence in sentences if sentence.strip())
        
        if not summary_text:
            summary_text = f"This document contains {word_count} words of content covering various topics and information"
        
        return f"{summary_text}. The document provides comprehensive information that could be valuable for understanding the subject matter."
    
    def _generate_fallback_strengths_weaknesses(self) -> Dict[str, List[str]]:
        """Generate fallback strengths and weaknesses."""
        return {
            "strengths": [
                "Document provides structured information on the topic",
                "Content appears to be well-organized and comprehensive",
                "Information is presented in a clear format"
            ],
            "weaknesses": [
                "Could benefit from more detailed analysis",
                "Some sections might need additional supporting evidence",
                "Document structure could be enhanced for better readability"
            ]
        }
    
    def _generate_fallback_questions(self) -> List[str]:
        """Generate fallback exploration questions."""
        return [
            "What are the key implications of the main findings presented?",
            "How could this information be applied in practical scenarios?",
            "What additional research or data would strengthen the conclusions?",
            "Are there alternative perspectives that should be considered?",
            "What are the potential long-term effects of implementing these ideas?"
        ]
    
    def _generate_fallback_recommendations(self) -> List[str]:
        """Generate fallback recommendations."""
        return [
            "Review and validate the key findings with additional sources",
            "Develop an implementation plan based on the document's insights",
            "Gather stakeholder feedback on the proposed approaches",
            "Create metrics to measure the effectiveness of recommendations",
            "Schedule regular reviews to assess progress and make adjustments"
        ]