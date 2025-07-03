from fastapi import APIRouter, HTTPException, UploadFile, File
from models.schemas import AnalyzeRequest, AnalyzeResponse, VideoMetadata
from services.transcriber import TranscriberService, VideoProcessingError
from services.viral import ViralAnalysisService
from services.gemini_utils import summarize_transcript, explain_why_viral, generate_content_idea
from utils import youtube
import logging
import tempfile
import os
from pathlib import Path
import PyPDF2
import docx
import pptx

router = APIRouter()
logger = logging.getLogger(__name__)

transcriber_service = TranscriberService()
viral_service = ViralAnalysisService()

@router.post("/analyze", response_model=AnalyzeResponse)
async def analyze_content(request: AnalyzeRequest):
    """Menganalisis konten YouTube."""
    if not request.youtube_url:
        raise HTTPException(status_code=400, detail="youtube_url must be provided")
    logger.info(f"Analyzing YouTube content: {request.youtube_url}")
    try:
        video_metadata = await youtube.get_video_metadata(request.youtube_url)
        if not video_metadata:
            raise HTTPException(status_code=404, detail="Invalid YouTube URL or video not found.")
        
        comments = await youtube.get_video_comments(video_metadata.video_id)
        transcript = await transcriber_service.get_transcript(request.youtube_url)
        
        overall_summary = await summarize_transcript(transcript)
        viral_explanation = await explain_why_viral(video_metadata.title, video_metadata.view_count or 0, video_metadata.like_count or 0, overall_summary)
        recommendations = await generate_content_idea("youtube", overall_summary, viral_explanation)
        
        # --- PERBAIKAN DI SINI ---
        # Mengubah cara pemanggilan fungsi agar sesuai dengan definisi baru
        viral_score = await viral_service.calculate_viral_score(
            content=transcript,
            metadata=video_metadata
        )
        
        if viral_score >= 80: viral_label = "Very High Potential"
        elif viral_score >= 60: viral_label = "Good Potential"
        else: viral_label = "Needs Improvement"
        
        return AnalyzeResponse(
            video_metadata=video_metadata, summary=overall_summary, timeline_summary=[],
            viral_score=viral_score, viral_label=viral_label, 
            viral_explanation=viral_explanation, recommendations=recommendations
        )
    except Exception as e:
        logger.error(f"An unexpected server error occurred: {e}", exc_info=True)
        # Memberikan detail error ke client untuk mempermudah debugging
        raise HTTPException(status_code=500, detail=f"An internal server error occurred: {type(e).__name__}")

@router.post("/analyze-document", response_model=AnalyzeResponse)
async def analyze_document(file: UploadFile = File(...)):
    """Menganalisis dokumen yang diunggah."""
    allowed_extensions = {'.pdf', '.doc', '.docx', '.txt', '.ppt', '.pptx'}
    file_extension = Path(file.filename or '').suffix.lower()
    if file_extension not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}")
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        try:
            document_text = await extract_text_from_document(temp_file_path, file_extension)
            if not document_text or len(document_text.strip()) < 20:
                raise HTTPException(status_code=422, detail="Document content is too short or empty.")
            
            overall_summary = await summarize_transcript(document_text)
            
            # Untuk dokumen, kita buat metadata dummy karena tidak relevan
            dummy_metadata = VideoMetadata(
                video_id="doc_analysis", title=file.filename or "Document", duration=0,
                thumbnail_url="", channel_name="", channel_id="", view_count=0, like_count=0,
                comment_count=0, subscriber_count=0, published_at=None, description=""
            )
            viral_score = await viral_service.calculate_viral_score(document_text, dummy_metadata)

            if viral_score >= 80: viral_label = "Very High Potential"
            elif viral_score >= 60: viral_label = "Good Potential"
            else: viral_label = "Needs Improvement"

            viral_explanation = "This document has strong potential to be repurposed into engaging digital content."
            recommendations = await generate_content_idea("document", overall_summary, viral_explanation)
            
            return AnalyzeResponse(
                summary=overall_summary, viral_score=viral_score, viral_label=viral_label,
                viral_explanation=viral_explanation, recommendations=recommendations, doc_summary=overall_summary
            )
        finally:
            os.unlink(temp_file_path)
    except Exception as e:
        logger.error(f"Document analysis failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Failed to analyze document.")

async def extract_text_from_document(file_path: str, file_extension: str) -> str:
    """Mengekstrak teks dari berbagai format dokumen."""
    text = ""
    try:
        if file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f: return f.read()
        elif file_extension == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages: text += page.extract_text() or ""
            return text
        elif file_extension in ['.doc', '.docx']:
            doc = docx.Document(file_path)
            for para in doc.paragraphs: text += para.text + "\n"
            return text
        elif file_extension in ['.ppt', '.pptx']:
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
            return text
        raise ValueError(f"Unsupported file type: {file_extension}")
    except Exception as e:
        raise Exception(f"Failed to process {file_extension} file: {e}")